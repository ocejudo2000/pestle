"""import openai
import streamlit as st

# Configura tu API key para OpenAI

openai_api_key = st.secrets["mykey"]

openai.api_key = openai_api_key




def generar_resumen_pestle(nombre_empresa, caracteres_resumen):
    dimensiones_pestle = ["Político", "Económico", "Sociocultural", "Tecnológico", "Legal", "Ecológico"]
    resumen_pestle = {}

    for dim in dimensiones_pestle:
        prompt = f"Genera un resumen en formato bullet point sobre el aspecto {dim} de la empresa {nombre_empresa} en {caracteres_resumen} caracteres sin final cortado, pon cada aspecto en una tabla. Adicional Muestra a sus principales competidores en una tabla:"
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=caracteres_resumen
        )
        resumen_pestle[dim] = response.choices[0].text.strip()

    return resumen_pestle

# Interfaz Streamlit
st.title("Análisis PESTLE")
nombre_empresa = st.text_input("Ingresa el nombre de la Empresa:")

caracteres_resumen = st.slider("Cantidad de caracteres para resumir por dimensión PESTLE:", 50, 500, 250)

if st.button("Generar Resumen PESTLE"):
    resumenes = generar_resumen_pestle(nombre_empresa, caracteres_resumen)
    for dim, resumen in resumenes.items():
        st.write(f"**{dim}**")
        st.write(resumen)

export_option = st.selectbox("¿Quieres exportar los resúmenes?", ["No", "Excel", "PowerPoint"])

if export_option == "Excel":
    st.write("Funcionalidad para exportar a Excel aún no implementada.")

elif export_option == "PowerPoint":
    st.write("Funcionalidad para exportar a PowerPoint aún no implementada.")
    """
"""
import openai
import streamlit as st
import pandas as pd
from pptx import Presentation

# Configura tu API key para OpenAI
openai_api_key = st.secrets["mykey"]
openai.api_key = openai_api_key

def generar_resumen_pestle(nombre_empresa, caracteres_resumen):
    dimensiones_pestle = ["Político", "Económico", "Sociocultural", "Tecnológico", "Legal", "Ecológico"]
    resumen_pestle = {}
    for dim in dimensiones_pestle:
        prompt = f"Genera un resumen en formato bullet point sobre el aspecto {dim} de la empresa {nombre_empresa} en {caracteres_resumen} caracteres  sin final cortado, pon cada aspecto en una tabla. Adicional Muestra a sus principales competidores en una tabla:"
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=caracteres_resumen
        )
        resumen_pestle[dim] = response.choices[0].text.strip()
    return resumen_pestle

# Interfaz Streamlit
st.title("Análisis PESTLE")
nombre_empresa = st.text_input("Ingresa el nombre de la Empresa:")
caracteres_resumen = st.slider("Cantidad de caracteres para resumir por dimensión PESTLE:", 50, 500, 250)

if st.button("Generar Resumen PESTLE"):
    resumenes = generar_resumen_pestle(nombre_empresa, caracteres_resumen)
    for dim, resumen in resumenes.items():
        st.write(f"**{dim}**")
        st.write(resumen)

export_option = st.selectbox("¿Quieres exportar los resúmenes?", ["No", "Excel", "PowerPoint"])

# ---------> HERE IS THE CHANGE <---------
if export_option == "Excel":
    # Create a Pandas DataFrame from the resumenes dict
    df = pd.DataFrame(list(resumenes.items()), columns=["Aspect", "Summary"])
    towrite = io.BytesIO()
    downloaded_file = df.to_excel(towrite, encoding='utf-8', index=False, header=True)
    towrite.seek(0)
    if st.download_button('Download Excel', towrite, file_name='resumenes_pestle.xlsx', mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'):
        pass

elif export_option == "PowerPoint":
    # Create a presentation object
    prs = Presentation()
    for dim, resumen in resumenes.items():
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = dim
        content.text = resumen
    
    # Save to BytesIO object
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    if st.download_button('Download PowerPoint', ppt_io, file_name='resumenes_pestle.pptx', mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'):
        pass

"""
import openai
import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import os

# Configura tu API key para OpenAI
openai_api_key = st.secrets["mykey"]
openai.api_key = openai_api_key

def generar_resumen_pestle(nombre_empresa, caracteres_resumen):
    dimensiones_pestle = ["Político", "Económico", "Sociocultural", "Tecnológico", "Legal", "Ecológico"]
    resumen_pestle = {}
    competidores = None  # Placeholder for now

    for dim in dimensiones_pestle:
        prompt = f"Genera un resumen en formato bullet point sobre el aspecto {dim} de la empresa {nombre_empresa} en {caracteres_resumen} caracteres. Además, dime quiénes son sus principales competidores:"
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=caracteres_resumen
        )
        resumen_pestle[dim] = response.choices[0].text.strip().split("\n")

    return resumen_pestle, competidores

# Interfaz Streamlit
st.title("Análisis PESTLE")
nombre_empresa = st.text_input("Ingresa el nombre de la Empresa:")
caracteres_resumen = st.slider("Cantidad de caracteres para resumir por dimensión PESTLE:", 50, 500, 250)

if st.button("Generar Resumen PESTLE"):
    resumenes, competidores = generar_resumen_pestle(nombre_empresa, caracteres_resumen)
    for dim, resumen in resumenes.items():
        st.write(f"**{dim}**")
        st.table(pd.DataFrame({"Resumen": resumen}))

    export_option = st.selectbox("¿Quieres exportar los resúmenes?", ["No", "Excel", "PowerPoint"])

    if export_option == "Excel":
        df = pd.DataFrame(list(resumenes.items()), columns=["Aspect", "Summary"])
        df.to_excel("resumenes_pestle.xlsx", index=False)

        with open("resumenes_pestle.xlsx", "rb") as f:
            st.download_button("Descargar Excel", f.read(), file_name="resumenes_pestle.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    elif export_option == "PowerPoint":
        prs = Presentation()
        for dim, resumen in resumenes.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = dim
            content = slide.shapes.placeholders[1]
            for point in resumen:
                p = content.text = content.text + "\n" + point

        prs.save("resumenes_pestle.pptx")
        
        with open("resumenes_pestle.pptx", "rb") as f:
            st.download_button("Descargar PowerPoint", f.read(), file_name="resumenes_pestle.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

